from PyPDF2 import PdfReader
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Cargar el PDF
pdf_path = "C:\\Users\\luis.alamo\\Downloads\\DELTAPRESENTACIONEMPRESA2025.pdf"
reader = PdfReader(pdf_path)

# Crear una presentación de PowerPoint
prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]  # Slide en blanco

# Estilo de texto
def add_textbox(slide, text):
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(6.5)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    frame = textbox.text_frame
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.font.size = Pt(18)
    p.font.name = 'Calibri'
    p.font.color.rgb = RGBColor(0, 0, 0)
    frame.text = text

# Añadir cada página del PDF como una diapositiva con texto
for page in reader.pages:
    slide = prs.slides.add_slide(blank_slide_layout)
    text = page.extract_text()
    if text:
        add_textbox(slide, text.strip())

# Guardar la presentación
pptx_path = "C:\\Users\\luis.alamo\\Desktop\\DELTAPRESENTACIONEMPRESA2025.pptx"
prs.save(pptx_path)

pptx_path